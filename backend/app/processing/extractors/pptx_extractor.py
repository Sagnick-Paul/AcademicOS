"""PPTX extractor using python-pptx."""
from __future__ import annotations

from pathlib import Path
from typing import Any, Dict
# pyrefly: ignore [missing-import]
from pptx import Presentation

from app.processing.base import BaseExtractor
from app.processing.exceptions import ExtractionFailed

class PPTXExtractor(BaseExtractor):
    """Extracts text, slide notes, and metadata from PowerPoint files."""

    async def extract(self, file_path: Path) -> Dict[str, Any]:
        try:
            prs = Presentation(file_path)
        except Exception as exc:
            raise ExtractionFailed(f"python-pptx failed to open file: {exc}") from exc

        try:
            pages = []
            full_text_list = []

            for i, slide in enumerate(prs.slides):
                slide_texts = []
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_texts.append(shape.text)
                
                slide_text = "\n".join(slide_texts)
                
                # Extract slide notes if present
                notes_text = ""
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text

                # Combine slide text and notes for the final page block
                full_text_list.append(slide_text)
                if notes_text:
                    full_text_list.append(f"[Notes]: {notes_text}")

                pages.append({
                    "page_index": i,
                    "text": slide_text,
                    "notes": notes_text,
                })

            full_text = "\n".join(full_text_list)
            
            # Simple title resolution (use first slide title or filename)
            title = None
            if prs.slides:
                slide = prs.slides[0]
                if slide.shapes.title:
                    title = slide.shapes.title.text

            return {
                "text": full_text,
                "metadata": {
                    "title": title or file_path.stem,
                    "author": None,
                    "pages": len(prs.slides),
                },
                "pages": pages,
            }
        except Exception as exc:
            raise ExtractionFailed(f"python-pptx failed during extraction: {exc}") from exc
