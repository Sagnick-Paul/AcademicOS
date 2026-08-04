"""Tests for Phase 2 — Intelligent Document Processing."""
from __future__ import annotations

import io
from pathlib import Path
from uuid import UUID

# pyrefly: ignore [missing-import]
import pytest
# pyrefly: ignore [missing-import]
import fitz
# pyrefly: ignore [missing-import]
from pptx import Presentation
# pyrefly: ignore [missing-import]
from httpx import AsyncClient

from app.processing.cleaners import clean_text
from app.processing.chunking import ChunkingService
from app.processing.exceptions import UnsupportedDocumentType, ExtractionFailed
from app.processing.extractors.text_extractor import TextExtractor
from app.processing.extractors.pymupdf_extractor import PyMuPDFExtractor
from app.processing.extractors.pptx_extractor import PPTXExtractor
from app.processing.processors.pdf_processor import PDFProcessor
from app.processing.dispatcher import DocumentDispatcher
from app.processing.pipeline import DocumentProcessingPipeline
from app.db.models.enums import DocumentUploadStatus


@pytest.fixture()
def temp_files_dir(tmp_path: Path) -> Path:
    """Fixture providing a temporary directory for generating test files."""
    return tmp_path


# ---------- Unit Tests: Cleaners & Chunking ----------


def test_clean_text() -> None:
    text = "  Hello \x00 world!  \n\n\nNew line\twith tab.\r\n\n\nThird line.  "
    cleaned, duration = clean_text(text)
    
    assert duration > 0.0
    assert "\x00" not in cleaned
    assert "\t" not in cleaned
    # 3+ consecutive newlines reduced to 2
    assert "\n\n\n" not in cleaned
    assert cleaned.startswith("Hello")
    assert cleaned.endswith("Third line.")


def test_recursive_chunking() -> None:
    chunker = ChunkingService(chunk_size=50, overlap=10)
    text = "This is a long piece of text that will be split recursively into multiple overlapping chunks."
    
    chunks, duration = chunker.chunk_text(text)
    assert len(chunks) > 1
    assert duration > 0.0
    for chunk in chunks:
        assert len(chunk.text) <= 50
        assert chunk.text != ""


# ---------- Unit Tests: Extractors ----------


async def test_text_extractor(temp_files_dir: Path) -> None:
    txt_path = temp_files_dir / "test.txt"
    content = "Hello plain text world!"
    txt_path.write_text(content, encoding="utf-8")

    extractor = TextExtractor()
    result = await extractor.extract(txt_path)
    
    assert result["text"] == content
    assert result["metadata"]["pages"] == 1
    assert len(result["pages"]) == 1


async def test_pdf_extractor(temp_files_dir: Path) -> None:
    pdf_path = temp_files_dir / "test.pdf"
    
    # Generate a real PDF in memory using fitz (PyMuPDF)
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((50, 50), "Hello world from PyMuPDF!")
    doc.set_metadata({"title": "Test Title", "author": "Test Author"})
    doc.save(pdf_path)
    doc.close()

    extractor = PyMuPDFExtractor()
    result = await extractor.extract(pdf_path)
    
    assert "Hello world from PyMuPDF" in result["text"]
    assert result["metadata"]["pages"] == 1
    assert result["metadata"]["title"] == "Test Title"
    assert result["metadata"]["author"] == "Test Author"


async def test_pptx_extractor(temp_files_dir: Path) -> None:
    pptx_path = temp_files_dir / "test.pptx"

    # Generate a real PPTX using python-pptx
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "Hello PPTX Title"
    
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "Slide notes content"
    
    prs.save(pptx_path)

    extractor = PPTXExtractor()
    result = await extractor.extract(pptx_path)
    
    assert "Hello PPTX Title" in result["text"]
    assert "Slide notes content" in result["text"]
    assert result["metadata"]["pages"] == 1


# ---------- Unit Tests: Pipeline & Dispatcher ----------


async def test_dispatcher_unsupported_type(temp_files_dir: Path) -> None:
    exe_path = temp_files_dir / "test.exe"
    exe_path.write_bytes(b"MZ\x90")

    dispatcher = DocumentDispatcher()
    with pytest.raises(UnsupportedDocumentType):
        await dispatcher.dispatch(exe_path, "exe", "test.exe", 3)


async def test_pipeline_error_wrapping(temp_files_dir: Path) -> None:
    non_existent = temp_files_dir / "does_not_exist.pdf"
    
    pipeline = DocumentProcessingPipeline()
    with pytest.raises(ExtractionFailed):
        await pipeline.run(non_existent, "pdf", "does_not_exist.pdf", 0)


# ---------- Integration: Upload + Background Processing ----------


async def register_and_login(client: AsyncClient, *, email: str) -> str:
    """Helper to register and login a user."""
    reg = await client.post(
        "/api/v1/auth/register",
        json={"email": email, "password": "S3cur3P@ss!", "full_name": "Test"},
    )
    assert reg.status_code == 201

    login = await client.post(
        "/api/v1/auth/login",
        json={"email": email, "password": "S3cur3P@ss!"},
    )
    assert login.status_code == 200
    return login.json()["access_token"]


async def test_upload_and_processing_pipeline(client: AsyncClient, upload_tmp_path: Path) -> None:
    """End-to-end integration test: upload, process in background, and save sidecar."""
    token = await register_and_login(client, email="pipeline@example.com")
    headers = {"Authorization": f"Bearer {token}"}

    # 1. Generate a real PDF to upload
    pdf_bytes = io.BytesIO()
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((50, 50), "End-to-end test document content.")
    doc.save(pdf_bytes)
    doc.close()
    pdf_bytes.seek(0)

    # 2. POST upload
    resp = await client.post(
        "/api/v1/documents/upload",
        files={"file": ("pipeline_test.pdf", pdf_bytes, "application/pdf")},
        headers=headers,
    )
    assert resp.status_code == 201, resp.text
    doc_id = resp.json()["id"]

    # 3. Fetch status — background tasks run synchronously in ASGI test client,
    # so the document should already be READY by the time we fetch it.
    get_resp = await client.get(f"/api/v1/documents/{doc_id}", headers=headers)
    assert get_resp.status_code == 200
    assert get_resp.json()["upload_status"] == DocumentUploadStatus.READY

    # 4. Check sidecar JSON file exists in local storage
    storage_path = get_resp.json()["storage_path"]
    sidecar_path = upload_tmp_path / f"{storage_path}.processing.json"
    assert sidecar_path.is_file(), f"Sidecar should exist: {sidecar_path}"
